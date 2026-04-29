"""rag.py

Core RAG logic — summarization + retrieval.

Features:
- Ingest PDF/PPTX into overlapping text chunks
- BM25 retrieval + local (NumPy) cosine-similarity embeddings via Ollama
- Optional summary generation (cached per chapter)
- Optional Graph RAG using NetworkX (cached per chapter)

This module intentionally keeps a simple, function-based API used by `backend/main.py`:
- startup(pdf_path)
- load_content(content_path, store_key=None)
- answer(question, mode=None, debug=False)
- teach(), teach_structured(), generate_quiz()
- graph_stats(), graph_export()

NOTE: This module uses module-level mutable state.
That is convenient for a single-user/local app, but is not multi-tenant safe.
"""

from __future__ import annotations

import json
import os
import re
import uuid
from pathlib import Path
from typing import Any

import numpy as np
import ollama
from pypdf import PdfReader
from rank_bm25 import BM25Okapi

try:
    import networkx as nx

    _NX_AVAILABLE = True
except ImportError:
    _NX_AVAILABLE = False
    nx = None  # type: ignore
    print("[WARN] networkx not installed — Graph RAG disabled.")

try:
    from pptx import Presentation
except Exception:  # pragma: no cover
    Presentation = None


# ── Config ─────────────────────────────────────────────────────────────

CHAT_MODEL = os.getenv("CHAT_MODEL", "mistral")
EMBED_MODEL = os.getenv("EMBED_MODEL", "nomic-embed-text")
VECTOR_STORE_ROOT = Path(os.getenv("VECTOR_STORE_DIR", ".rag_store"))

CHUNK_CHARS = int(os.getenv("RAG_CHUNK_CHARS", "1200"))
CHUNK_OVERLAP = int(os.getenv("RAG_CHUNK_OVERLAP", "200"))

ENABLE_SUMMARY = os.getenv("RAG_ENABLE_SUMMARY", "1").strip().lower() not in {"0", "false", "no"}
ENABLE_GRAPH_RAG = os.getenv("RAG_ENABLE_GRAPH", "1").strip().lower() not in {"0", "false", "no"}

MIN_VECTOR_SCORE = float(os.getenv("RAG_MIN_VECTOR_SCORE", "0.18"))
MIN_BM25_SCORE = float(os.getenv("RAG_MIN_BM25_SCORE", "0.05"))

MODE_INSTRUCTIONS: dict[str, str] = {
    "default": "",
    "explain7": "Explain like I'm 7 years old. Use very simple words and short sentences.",
    "example": "Include one simple real-life example that matches the textbook context.",
    "stepbystep": "Explain step by step using a numbered list.",
}


# ── State (reloaded per chapter/content) ───────────────────────────────

texts: list[str] = []
metadatas: list[dict[str, Any]] = []
ids: list[str] = []

bm25: BM25Okapi | None = None
vector_matrix: np.ndarray | None = None

chunk_summaries: list[str] = []
final_summary: str = ""

knowledge_graph: Any = None  # nx.DiGraph
entity_to_chunks: dict[str, list[int]] = {}

CURRENT_CONTENT_PATH: str | None = None
CURRENT_STORE_KEY: str | None = None

# Cache file paths (set by set_store_dir)
_STORE_DIR: Path = VECTOR_STORE_ROOT
_MANIFEST_PATH: Path = VECTOR_STORE_ROOT / "manifest.json"
_EMBEDDINGS_PATH: Path = VECTOR_STORE_ROOT / "embeddings.npy"
_SUMMARY_PATH: Path = VECTOR_STORE_ROOT / "summary.json"
_GRAPH_PATH: Path = VECTOR_STORE_ROOT / "knowledge_graph.json"


def set_store_dir(store_dir: Path) -> None:
    global _STORE_DIR, _MANIFEST_PATH, _EMBEDDINGS_PATH, _SUMMARY_PATH, _GRAPH_PATH
    _STORE_DIR = Path(store_dir)
    _MANIFEST_PATH = _STORE_DIR / "manifest.json"
    _EMBEDDINGS_PATH = _STORE_DIR / "embeddings.npy"
    _SUMMARY_PATH = _STORE_DIR / "summary.json"
    _GRAPH_PATH = _STORE_DIR / "knowledge_graph.json"


# ── Ingest ─────────────────────────────────────────────────────────────


def _clean_text(text: str) -> str:
    text = text or ""
    text = text.replace("\u00a0", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _chunk_text(text: str, chunk_chars: int = CHUNK_CHARS, overlap: int = CHUNK_OVERLAP) -> list[str]:
    text = _clean_text(text)
    if not text:
        return []
    if chunk_chars <= 0:
        return [text]

    chunks: list[str] = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + chunk_chars, n)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= n:
            break
        start = max(0, end - overlap)
    return chunks


def _discover_files(content_path: str) -> list[str]:
    p = Path(content_path).expanduser().resolve()
    if p.is_file():
        return [str(p)]
    if p.is_dir():
        files: list[str] = []
        for pat in ("*.pdf", "*.pptx"):
            files.extend(str(x) for x in p.rglob(pat))
        return sorted(set(files))
    raise FileNotFoundError(f"Content path not found: {content_path}")


def _extract_pdf_units(path: str) -> list[tuple[str, dict[str, Any]]]:
    reader = PdfReader(path)
    units: list[tuple[str, dict[str, Any]]] = []
    for idx, page in enumerate(reader.pages, start=1):
        t = _clean_text(page.extract_text() or "")
        if not t:
            continue
        units.append((t, {"source": os.path.basename(path), "path": path, "page": idx}))
    return units


def _extract_pptx_units(path: str) -> list[tuple[str, dict[str, Any]]]:
    if Presentation is None:
        raise RuntimeError("python-pptx not installed. Add it to requirements.txt")

    prs = Presentation(path)
    units: list[tuple[str, dict[str, Any]]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                txt = _clean_text(shape.text_frame.text or "")
                if txt:
                    parts.append(txt)
        t = _clean_text("\n".join(parts))
        if not t:
            continue
        units.append((t, {"source": os.path.basename(path), "path": path, "slide": idx}))
    return units


def ingest_paths(file_paths: list[str], *, reset: bool = True, rebuild_summary: bool = True) -> None:
    """Ingest a list of local file paths into memory + indexes."""
    global texts, metadatas, ids

    if reset:
        texts = []
        metadatas = []
        ids = []

    extracted_units = 0
    for fp in file_paths:
        ext = Path(fp).suffix.lower()
        if ext == ".pdf":
            units = _extract_pdf_units(fp)
        elif ext == ".pptx":
            units = _extract_pptx_units(fp)
        else:
            continue

        extracted_units += len(units)

        for unit_text, meta in units:
            for chunk in _chunk_text(unit_text):
                texts.append(chunk)
                metadatas.append(meta)
                ids.append(str(uuid.uuid4()))

    if not texts:
        raise ValueError("No text extracted from provided documents. (Are they scanned/image-only?)")

    print(f"[RAG] Extracted {extracted_units} text units -> {len(texts)} chunks")

    build_index()
    build_vector_index()

    if ENABLE_SUMMARY:
        if not load_summary():
            if rebuild_summary:
                build_summary()
                save_summary()

    if ENABLE_GRAPH_RAG and _NX_AVAILABLE:
        if not load_knowledge_graph():
            build_knowledge_graph()


def ingest(content_path: str, *, reset: bool = True, rebuild_summary: bool = True) -> None:
    paths = _discover_files(content_path)
    if not paths:
        raise ValueError(f"No .pdf/.pptx files found under: {content_path}")
    ingest_paths(paths, reset=reset, rebuild_summary=rebuild_summary)


# ── Indexes ────────────────────────────────────────────────────────────


def build_index() -> None:
    global bm25
    tokenized = [t.lower().split() for t in texts]
    bm25 = BM25Okapi(tokenized)
    print("[RAG] BM25 index built.")


def build_vector_index(*, batch_size: int = 48) -> None:
    """Build or load a local vector index (cosine similarity) persisted to _STORE_DIR."""
    global vector_matrix

    if not texts:
        vector_matrix = None
        return

    # Load cached embeddings if they match.
    try:
        if _MANIFEST_PATH.exists() and _EMBEDDINGS_PATH.exists():
            manifest = json.loads(_MANIFEST_PATH.read_text(encoding="utf-8"))
            if manifest.get("embed_model") == EMBED_MODEL and int(manifest.get("chunks") or 0) == len(texts):
                vector_matrix = np.load(_EMBEDDINGS_PATH)
                print("[RAG] Loaded cached local vector index.")
                return
    except Exception:
        pass

    print(f"[RAG] Building local vector index with {len(texts)} chunks...")

    all_embeddings: list[list[float]] = []
    for start in range(0, len(texts), batch_size):
        batch_docs = texts[start : start + batch_size]
        emb_res = ollama.embed(model=EMBED_MODEL, input=batch_docs)
        batch_embeddings = emb_res.get("embeddings") if isinstance(emb_res, dict) else getattr(emb_res, "embeddings", None)
        if not batch_embeddings:
            raise RuntimeError("No embeddings returned from Ollama. Is the embedding model available?")
        all_embeddings.extend(batch_embeddings)

    mat = np.asarray(all_embeddings, dtype=np.float32)
    norms = np.linalg.norm(mat, axis=1, keepdims=True)
    mat = mat / np.clip(norms, 1e-12, None)
    vector_matrix = mat

    _STORE_DIR.mkdir(parents=True, exist_ok=True)
    np.save(_EMBEDDINGS_PATH, vector_matrix)
    _MANIFEST_PATH.write_text(
        json.dumps({"embed_model": EMBED_MODEL, "chunks": len(texts)}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print("[RAG] Local vector index built.")


def _format_meta(meta: dict[str, Any] | None) -> str:
    meta = meta or {}
    src = meta.get("source") or meta.get("path") or "unknown"
    if meta.get("page"):
        return f"{src} (page {meta['page']})"
    if meta.get("slide"):
        return f"{src} (slide {meta['slide']})"
    return str(src)


def retrieve_vector_scored(query: str, k: int = 3) -> list[dict[str, Any]]:
    if vector_matrix is None or not texts:
        return []

    try:
        emb_res = ollama.embed(model=EMBED_MODEL, input=query)
        query_embs = emb_res.get("embeddings") if isinstance(emb_res, dict) else getattr(emb_res, "embeddings", None)
        if not query_embs:
            return []
        q = np.asarray(query_embs[0], dtype=np.float32)
        q = q / max(float(np.linalg.norm(q)), 1e-12)

        scores = vector_matrix @ q
        kk = max(1, int(k))
        top_idx = np.argsort(-scores)[:kk]

        out: list[dict[str, Any]] = []
        for i in top_idx:
            idx = int(i)
            out.append(
                {
                    "text": texts[idx],
                    "meta": metadatas[idx] if idx < len(metadatas) else {},
                    "score": float(scores[idx]),
                    "retriever": "vector",
                }
            )
        return out
    except Exception:
        return []


def retrieve_bm25_scored(query: str, k: int = 3) -> list[dict[str, Any]]:
    if bm25 is None:
        return []

    tokenized_query = query.lower().split()
    scores = bm25.get_scores(tokenized_query)
    kk = max(1, int(k))
    top_idx = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)[:kk]

    out: list[dict[str, Any]] = []
    for idx in top_idx:
        out.append(
            {
                "text": texts[idx],
                "meta": metadatas[idx] if idx < len(metadatas) else {},
                "score": float(scores[idx]),
                "retriever": "bm25",
            }
        )
    return out


def _good_enough_hits(hits: list[dict[str, Any]]) -> bool:
    if not hits:
        return False
    retriever = hits[0].get("retriever")
    best = float(hits[0].get("score") or 0.0)
    if retriever == "vector":
        return best >= MIN_VECTOR_SCORE
    if retriever == "bm25":
        return best >= MIN_BM25_SCORE
    return False


# ── Summary (cached) ───────────────────────────────────────────────────


def build_summary(chunk_size: int = 5) -> None:
    global chunk_summaries, final_summary

    chunks = ["\n".join(texts[i : i + chunk_size]) for i in range(0, len(texts), chunk_size)]
    print(f"[RAG] Summarizing {len(chunks)} chunks...")

    def summarize_chunk(chunk_text: str, idx: int) -> str:
        prompt = f"""You are summarizing a CBSE primary school textbook chapter.
Summarize in 3-4 simple sentences. Focus on key concepts only.

Text:
{chunk_text[:3000]}

Summary:"""
        res = ollama.chat(model=CHAT_MODEL, messages=[{"role": "user", "content": prompt}])
        summary = res["message"]["content"].strip()
        print(f"  ✓ Chunk {idx + 1} summarized")
        return summary

    chunk_summaries = [summarize_chunk(c, i) for i, c in enumerate(chunks)]

    summaries_for_final = chunk_summaries
    if len(chunk_summaries) > 10:
        print("[RAG] Large content — doing mid-level merge...")
        mid_chunks = ["\n".join(chunk_summaries[i : i + 5]) for i in range(0, len(chunk_summaries), 5)]
        mid_summaries: list[str] = []
        for idx, mc in enumerate(mid_chunks):
            res = ollama.chat(
                model=CHAT_MODEL,
                messages=[{"role": "user", "content": f"Summarize in 2-3 sentences:\n\n{mc}\n\nSummary:"}],
            )
            mid_summaries.append(res["message"]["content"].strip())
            print(f"    ✓ Mid-merge {idx + 1}/{len(mid_chunks)}")
        summaries_for_final = mid_summaries

    combined = "\n\n".join(f"Section {i + 1}: {s}" for i, s in enumerate(summaries_for_final))
    prompt = f"""You are a CBSE primary school teacher.
Write a single child-friendly chapter overview in 5-6 sentences.

Section Summaries:
{combined}

Final Chapter Overview:"""
    res = ollama.chat(model=CHAT_MODEL, messages=[{"role": "user", "content": prompt}])
    final_summary = res["message"]["content"].strip()
    print("[RAG] Summarization complete!")


def save_summary() -> None:
    try:
        _STORE_DIR.mkdir(parents=True, exist_ok=True)
        _SUMMARY_PATH.write_text(
            json.dumps(
                {
                    "chunks": len(texts),
                    "final_summary": final_summary,
                    "chunk_summaries": chunk_summaries,
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )
    except Exception as e:
        print(f"[WARN] Failed to save summary cache: {e}")


def load_summary() -> bool:
    global final_summary, chunk_summaries

    if not _SUMMARY_PATH.exists():
        return False

    try:
        data = json.loads(_SUMMARY_PATH.read_text(encoding="utf-8"))
        if int(data.get("chunks") or 0) != len(texts):
            return False
        fs = (data.get("final_summary") or "").strip()
        cs = data.get("chunk_summaries")
        if not isinstance(cs, list):
            cs = []
        chunk_summaries = [str(s) for s in cs]
        final_summary = fs
        if final_summary:
            print("[RAG] Loaded cached summary.")
            return True
        return False
    except Exception as e:
        print(f"[WARN] Failed to load summary cache: {e}")
        return False


# ── Graph RAG (cached) ────────────────────────────────────────────────


def _extract_triples_from_text(text: str) -> list[tuple[str, str, str]]:
    prompt = f"""Extract subject-relation-object triples from this educational text.
Return ONLY a JSON array like: [{{"e1":"photosynthesis","rel":"produces","e2":"oxygen"}}]
Max 6 triples. Keep entities short (1-3 words). No markdown.

Text:
{text[:900]}

Triples:"""

    try:
        res = ollama.chat(model=CHAT_MODEL, messages=[{"role": "user", "content": prompt}])
        raw = res["message"]["content"].strip()
        m = re.search(r"\[.*?\]", raw, re.DOTALL)
        if not m:
            return []
        triples = json.loads(m.group())
        if not isinstance(triples, list):
            return []
        out: list[tuple[str, str, str]] = []
        for t in triples:
            if not isinstance(t, dict):
                continue
            e1 = str(t.get("e1", "")).strip().lower()
            rel = str(t.get("rel", "")).strip().lower()
            e2 = str(t.get("e2", "")).strip().lower()
            if e1 and rel and e2:
                out.append((e1, rel, e2))
        return out
    except Exception as e:
        print(f"[Graph RAG] Triple extraction error: {e}")
        return []


def build_knowledge_graph(batch_size: int = 6) -> None:
    global knowledge_graph, entity_to_chunks

    if not _NX_AVAILABLE or not texts:
        return

    print(f"[Graph RAG] Building knowledge graph from {len(texts)} chunks (batch={batch_size})...")
    G = nx.DiGraph()  # type: ignore
    e2c: dict[str, list[int]] = {}

    for start in range(0, len(texts), batch_size):
        batch_idx = list(range(start, min(start + batch_size, len(texts))))
        combined = "\n\n".join(texts[i] for i in batch_idx)
        triples = _extract_triples_from_text(combined)

        for e1, rel, e2 in triples:
            for entity in (e1, e2):
                if entity not in G:
                    G.add_node(entity)
                if entity not in e2c:
                    e2c[entity] = []
                for idx in batch_idx:
                    if idx not in e2c[entity]:
                        e2c[entity].append(idx)
            G.add_edge(e1, e2, relation=rel)

        done = min(start + batch_size, len(texts))
        print(f"  [Graph RAG] {done}/{len(texts)} chunks processed")

    knowledge_graph = G
    entity_to_chunks = e2c
    print(f"[Graph RAG] ✓ Graph: {G.number_of_nodes()} nodes, {G.number_of_edges()} edges")
    save_knowledge_graph()


def save_knowledge_graph() -> None:
    if not _NX_AVAILABLE or knowledge_graph is None:
        return
    try:
        _STORE_DIR.mkdir(parents=True, exist_ok=True)
        data = {
            "nodes": list(knowledge_graph.nodes()),
            "edges": [
                {"e1": u, "rel": d.get("relation", ""), "e2": v}
                for u, v, d in knowledge_graph.edges(data=True)
            ],
            "entity_to_chunks": entity_to_chunks,
        }
        _GRAPH_PATH.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"[Graph RAG] Graph saved → {_GRAPH_PATH}")
    except Exception as e:
        print(f"[WARN] Failed to save graph cache: {e}")


def load_knowledge_graph() -> bool:
    global knowledge_graph, entity_to_chunks

    if not _NX_AVAILABLE or not _GRAPH_PATH.exists():
        return False

    try:
        data = json.loads(_GRAPH_PATH.read_text(encoding="utf-8"))
        G = nx.DiGraph()  # type: ignore
        G.add_nodes_from(data.get("nodes", []))
        for edge in data.get("edges", []):
            if not isinstance(edge, dict):
                continue
            G.add_edge(edge.get("e1"), edge.get("e2"), relation=edge.get("rel", ""))
        knowledge_graph = G
        etc = data.get("entity_to_chunks", {})
        entity_to_chunks = {str(k): list(v) for k, v in etc.items()} if isinstance(etc, dict) else {}
        print(f"[Graph RAG] Loaded graph: {G.number_of_nodes()} nodes, {G.number_of_edges()} edges")
        return True
    except Exception as e:
        print(f"[WARN] Failed to load graph cache: {e}")
        return False


def retrieve_graph(query: str, k: int = 3, hops: int = 2) -> list[str]:
    if not _NX_AVAILABLE or knowledge_graph is None or not entity_to_chunks:
        return []

    prompt = f"""List the key concepts/entities in this question as a JSON array.
Max 5 items, each 1-3 words, lowercase. No markdown.
Question: {query}
Entities:"""

    query_entities: list[str] = []
    try:
        res = ollama.chat(model=CHAT_MODEL, messages=[{"role": "user", "content": prompt}])
        raw = res["message"]["content"].strip()
        m = re.search(r"\[.*?\]", raw, re.DOTALL)
        if m:
            query_entities = [str(e).strip().lower() for e in json.loads(m.group()) if e]
    except Exception:
        pass

    if not query_entities:
        query_entities = [w.lower() for w in re.findall(r"\b\w{4,}\b", query)][:5]

    all_nodes = set(knowledge_graph.nodes())
    seed_nodes: set[str] = set()
    for qe in query_entities:
        if qe in all_nodes:
            seed_nodes.add(qe)
        else:
            for node in all_nodes:
                if qe in node or node in qe:
                    seed_nodes.add(node)

    if not seed_nodes:
        return []

    visited: set[str] = set()
    frontier = set(seed_nodes)
    for _ in range(hops):
        nxt: set[str] = set()
        for node in frontier:
            if node not in visited:
                visited.add(node)
                nxt.update(knowledge_graph.successors(node))
                nxt.update(knowledge_graph.predecessors(node))
        frontier = nxt - visited
    visited.update(frontier)

    chunk_indices: set[int] = set()
    for node in visited:
        for idx in entity_to_chunks.get(node, []):
            if 0 <= int(idx) < len(texts):
                chunk_indices.add(int(idx))

    retrieved = [texts[i] for i in sorted(chunk_indices)[: max(1, int(k))]]
    print(f"[Graph RAG] Query entities: {query_entities} → {len(visited)} nodes → {len(retrieved)} chunks")
    return retrieved


def graph_stats() -> dict:
    if not _NX_AVAILABLE or knowledge_graph is None:
        return {"available": False, "nodes": 0, "edges": 0, "entities": []}
    return {
        "available": True,
        "nodes": int(knowledge_graph.number_of_nodes()),
        "edges": int(knowledge_graph.number_of_edges()),
        "entities": sorted(knowledge_graph.nodes())[:30],
    }


def graph_export(*, max_nodes: int = 250, max_edges: int = 800) -> dict[str, Any]:
    if not _NX_AVAILABLE or knowledge_graph is None:
        return {"available": False, "nodes": [], "links": [], "stats": {"nodes": 0, "edges": 0}}

    G = knowledge_graph
    total_nodes = int(G.number_of_nodes())
    total_edges = int(G.number_of_edges())

    try:
        max_nodes_i = max(1, int(max_nodes))
    except Exception:
        max_nodes_i = 250
    try:
        max_edges_i = max(1, int(max_edges))
    except Exception:
        max_edges_i = 800

    nodes_sorted = sorted(G.nodes(), key=lambda n: (G.degree(n), str(n)), reverse=True)
    kept_nodes = set(nodes_sorted[:max_nodes_i])

    nodes_out = [{"id": n} for n in kept_nodes]

    links_out: list[dict[str, Any]] = []
    for u, v, d in G.edges(data=True):
        if u in kept_nodes and v in kept_nodes:
            links_out.append({"source": u, "target": v, "relation": d.get("relation", "")})
            if len(links_out) >= max_edges_i:
                break

    return {
        "available": True,
        "nodes": nodes_out,
        "links": links_out,
        "stats": {
            "nodes": total_nodes,
            "edges": total_edges,
            "exported_nodes": len(nodes_out),
            "exported_edges": len(links_out),
        },
    }


# ── Answer ─────────────────────────────────────────────────────────────


def answer(question: str, mode: str | None = None, *, debug: bool = False) -> dict[str, Any]:
    hits = retrieve_vector_scored(question, k=4)
    if hits and not _good_enough_hits(hits):
        hits = []
    if not hits:
        hits = retrieve_bm25_scored(question, k=4)
        if hits and not _good_enough_hits(hits):
            hits = []

    sources = [
        {
            "source": _format_meta(h.get("meta")),
            "score": float(h.get("score") or 0.0),
            "retriever": h.get("retriever"),
            "snippet": (h.get("text") or "")[:400],
        }
        for h in hits
    ]

    graph_chunks: list[str] = []
    if ENABLE_GRAPH_RAG:
        try:
            graph_chunks = retrieve_graph(question, k=3)
        except Exception as e:
            print(f"[WARN] Graph retrieval failed: {e}")

    if not hits and not graph_chunks:
        return {
            "answer": "I couldn't find this in the uploaded material. Try rephrasing your question or uploading the specific page/slide that mentions it.",
            "sources": sources if debug else [],
            "graph_used": False,
            "graph_chunks_count": 0,
        }

    context_blocks: list[str] = []
    seen: set[str] = set()

    for h in hits:
        meta_label = _format_meta(h.get("meta"))
        t = (h.get("text") or "").strip()
        if t and t not in seen:
            seen.add(t)
            context_blocks.append(f"[Source: {meta_label}]\n{t}")

    for gc in graph_chunks:
        t = (gc or "").strip()
        if t and t not in seen:
            seen.add(t)
            context_blocks.append(f"[Graph context]\n{t}")

    context = "\n\n".join(context_blocks)

    mode_key = (mode or "default").strip().lower()
    mode_instruction = MODE_INSTRUCTIONS.get(mode_key, "")

    prompt = f"""You are a helpful tutor.

Rules:
- Answer ONLY using the provided Document context.
- If the answer is not in the Document context, say you couldn't find it in the uploaded material.
- Do not use outside knowledge.
- Use short paragraphs or bullet points where it helps clarity.
- Never use filler phrases like "Great question!", "Lets embark", "In this lesson we will".
- If asked to repeat or explain again, give a cleaner shorter version.
- Keep answers under 150 words unless the question genuinely needs more detail.
{('\nMode:\n- ' + mode_instruction) if mode_instruction else ''}

Document context:
{context}

Question: {question}

Answer:"""

    res = ollama.chat(model=CHAT_MODEL, messages=[{"role": "user", "content": prompt}])
    ans = res["message"]["content"].strip()

    return {
        "answer": ans,
        "sources": sources if debug else [],
        "graph_used": len(graph_chunks) > 0,
        "graph_chunks_count": len(graph_chunks),
    }


# ── Teach / Quiz ───────────────────────────────────────────────────────


def teach() -> str:
    if not final_summary:
        raise ValueError("Summary not ready")

    chunk_block = "\n".join(f"- {s}" for s in (chunk_summaries[:12] if chunk_summaries else []))
    prompt = f"""You are a friendly CBSE tutor for Classes 1–5.

Create a short lesson titled "Today's Lesson" using ONLY the provided summaries.

Requirements:
- Keep it child-friendly and curriculum-grounded.
- Use headings and short bullet points.
- Include: (1) Key ideas, (2) Simple examples, (3) 3 quick check questions.
- Avoid scary, political, or adult content.
- Do not mention that you are an AI.
- Keep it under ~250 words.

Final Summary:
{final_summary}

Section Summaries:
{chunk_block}

Lesson:"""

    res = ollama.chat(model=CHAT_MODEL, messages=[{"role": "user", "content": prompt}])
    return res["message"]["content"].strip()


def teach_structured() -> dict:
    if not final_summary:
        raise ValueError("Summary not ready")

    chunk_block = "\n".join(f"- {s}" for s in (chunk_summaries[:10] if chunk_summaries else []))

    prompt = f"""You are a friendly CBSE tutor for Classes 1-5.

Create a structured lesson from ONLY the provided summaries.

Return STRICT JSON ONLY (no markdown, no explanation) in this exact shape:
{{
  "title": "Today's Lesson: <topic name>",
  "intro": "A 2-3 sentence friendly introduction to hook the student.",
  "key_concepts": [
    {{"term": "concept name", "meaning": "simple 1-sentence explanation"}},
    {{"term": "concept name", "meaning": "simple 1-sentence explanation"}},
    {{"term": "concept name", "meaning": "simple 1-sentence explanation"}}
  ],
  "examples": [
    {{"icon": "🌱", "title": "Example title", "body": "A short real-life example."}},
    {{"icon": "💧", "title": "Example title", "body": "A short real-life example."}}
  ],
  "fun_fact": "One surprising, child-friendly fun fact from the chapter.",
  "summary_points": ["Point 1", "Point 2", "Point 3", "Point 4"],
  "check_questions": [
    {{"q": "Question?", "a": "Short answer."}},
    {{"q": "Question?", "a": "Short answer."}},
    {{"q": "Question?", "a": "Short answer."}}
  ]
}}

Rules:
- Use ONLY facts from the summaries below.
- Child-friendly language (Classes 1-5).
- No adult, political, or violent content.
- Keep each field concise.

Final Summary:
{final_summary}

Section Summaries:
{chunk_block}
"""

    last_error: Exception | None = None
    for attempt in range(2):
        try:
            if attempt == 0:
                res = ollama.chat(model=CHAT_MODEL, messages=[{"role": "user", "content": prompt}], format="json")
            else:
                res = ollama.chat(model=CHAT_MODEL, messages=[{"role": "user", "content": prompt}])

            raw = (res.get("message", {}) or {}).get("content", "").strip()
            try:
                data = json.loads(raw)
            except Exception:
                data = _extract_json(raw)

            if not isinstance(data, dict):
                raise ValueError("Model did not return a JSON object")

            data.setdefault("title", "Today's Lesson")
            data.setdefault("intro", "")
            data.setdefault("key_concepts", [])
            data.setdefault("examples", [])
            data.setdefault("fun_fact", "")
            data.setdefault("summary_points", [])
            data.setdefault("check_questions", [])
            return data
        except Exception as e:
            last_error = e
            print(f"[WARN] teach_structured attempt {attempt + 1} failed: {e}")

    raise RuntimeError(f"Structured lesson generation failed: {last_error}")


def _extract_json(text: str) -> dict:
    m = re.search(r"```(?:json)?\s*(\{.*\})\s*```", text, flags=re.DOTALL | re.IGNORECASE)
    if m:
        return json.loads(m.group(1))

    start = text.find("{")
    end = text.rfind("}")
    if start != -1 and end != -1 and end > start:
        return json.loads(text[start : end + 1])

    raise ValueError("No JSON found in model output")


def generate_quiz() -> dict:
    if not final_summary:
        raise ValueError("Summary not ready")

    prompt = f"""You are a CBSE tutor for Classes 1–5.

Generate a short quiz based ONLY on the chapter summary.

Return STRICT JSON ONLY (no markdown, no explanation) with this exact shape:
{{
  "mcq": [{{"question": "...", "options": ["A","B","C","D"], "answerIndex": 0, "explanation": "..."}}],
  "trueFalse": [{{"statement": "...", "answer": true, "explanation": "..."}}],
  "fillBlanks": [{{"prompt": "... ___ ...", "answer": "..."}}]
}}

Rules:
- 5 MCQs, 4 True/False, 4 Fill-in-the-blanks.
- Keep language simple.
- No adult, violent, political, or unsafe content.
- Make options plausible but only one correct.

Chapter Summary:
{final_summary}
"""

    last_error: Exception | None = None
    for attempt in range(2):
        try:
            if attempt == 0:
                res = ollama.chat(model=CHAT_MODEL, messages=[{"role": "user", "content": prompt}], format="json")
            else:
                res = ollama.chat(model=CHAT_MODEL, messages=[{"role": "user", "content": prompt}])

            raw = (res.get("message", {}) or {}).get("content", "").strip()
            try:
                quiz = json.loads(raw)
            except Exception:
                quiz = _extract_json(raw)

            return _normalize_quiz(quiz)
        except Exception as e:
            last_error = e

    raise RuntimeError(f"Quiz generation failed: {last_error}")


def _normalize_quiz(quiz: dict) -> dict:
    quiz = quiz if isinstance(quiz, dict) else {}
    quiz.setdefault("mcq", [])
    quiz.setdefault("trueFalse", [])
    quiz.setdefault("fillBlanks", [])

    if not isinstance(quiz["mcq"], list):
        quiz["mcq"] = []
    if not isinstance(quiz["trueFalse"], list):
        quiz["trueFalse"] = []
    if not isinstance(quiz["fillBlanks"], list):
        quiz["fillBlanks"] = []

    for q in quiz["mcq"]:
        if not isinstance(q, dict):
            continue
        if not isinstance(q.get("options"), list):
            q["options"] = []
        q["options"] = [str(o) for o in q["options"]][:4]
        try:
            q["answerIndex"] = int(q.get("answerIndex", 0))
        except Exception:
            q["answerIndex"] = 0
        if q["answerIndex"] < 0 or q["answerIndex"] >= max(len(q["options"]), 1):
            q["answerIndex"] = 0

    for q in quiz["trueFalse"]:
        if not isinstance(q, dict):
            continue
        q["answer"] = bool(q.get("answer", False))

    for q in quiz["fillBlanks"]:
        if not isinstance(q, dict):
            continue
        q["answer"] = str(q.get("answer", "")).strip()

    return quiz


# ── Content switching ───────────────────────────────────────────────────


def load_content(content_path: str, *, store_key: str | None = None) -> None:
    """Load a chapter/content set into memory.

    Uses per-content caches under VECTOR_STORE_ROOT/<store_key>.
    """
    global CURRENT_CONTENT_PATH, CURRENT_STORE_KEY

    content_path = str(content_path or "").strip()
    if not content_path:
        raise ValueError("content_path is required")

    resolved = str(Path(content_path).expanduser().resolve())
    if CURRENT_CONTENT_PATH and str(Path(CURRENT_CONTENT_PATH).expanduser().resolve()) == resolved:
        return

    if store_key is None:
        p = Path(content_path)
        store_key = p.stem if p.suffix else p.name

    safe_key = re.sub(r"[^a-zA-Z0-9_\-]+", "_", str(store_key)).strip("_") or "default"
    store_dir = VECTOR_STORE_ROOT / safe_key
    set_store_dir(store_dir)

    summary_cached = _SUMMARY_PATH.exists()

    print(f"[RAG] Loading content: {content_path} (cache={safe_key})")
    ingest(content_path, reset=True, rebuild_summary=not summary_cached)

    CURRENT_CONTENT_PATH = resolved
    CURRENT_STORE_KEY = safe_key


def startup(pdf_path: str) -> None:
    # Backwards-compatible: can be file or folder.
    load_content(pdf_path)
